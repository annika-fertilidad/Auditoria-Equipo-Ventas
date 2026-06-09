#!/usr/bin/env python3
"""Genera la presentación PXI con identidad de marca Fertilidad Integral."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Brand palette ──
CREAM   = RGBColor(0xF2, 0xF3, 0xE9)
IVORY   = RGBColor(0xF4, 0xF3, 0xE1)
FOREST  = RGBColor(0x20, 0x28, 0x1B)
SAGE    = RGBColor(0x73, 0x8D, 0x84)
SAGEL   = RGBColor(0x8F, 0xA8, 0x9F)
FSOFT   = RGBColor(0x3D, 0x4E, 0x36)
BORDER  = RGBColor(0xE8, 0xE7, 0xD8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OK      = RGBColor(0x6B, 0x9E, 0x6E)
WARN    = RGBColor(0xC9, 0xA2, 0x4B)
BAD     = RGBColor(0xC0, 0x57, 0x4A)
CDMX    = RGBColor(0x4A, 0x7B, 0x9D)
FONT    = "DM Sans"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h, text, size=16, color=FOREST, bold=False, italic=False,
        align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.name = font; f.bold = bold; f.italic = italic
        f.color.rgb = color
    return tb


def rect(s, x, y, w, h, fill, line=None, line_w=None, radius=True):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    if radius:
        try: shp.adjustments[0] = 0.06
        except Exception: pass
    return shp


def tag(s, x, y, text):
    t = box(s, x, y, Inches(4), Inches(0.35), text.upper(), size=11, color=SAGE, bold=True)
    t.text_frame.paragraphs[0].runs[0].font.name = FONT
    return t

# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════
s = slide(FOREST)
rect(s, Inches(0.9), Inches(1.0), Inches(0.9), Inches(0.9), SAGE)
box(s, Inches(0.9), Inches(1.0), Inches(0.9), Inches(0.9), "FI", size=26, color=CREAM,
    bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(2.0), Inches(1.15), Inches(8), Inches(0.6), "FERTILIDAD INTEGRAL",
    size=14, color=SAGEL, bold=True)
box(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.8),
    "Índice de Experiencia\ndel Paciente (PXI)", size=52, color=CREAM, bold=True, line_spacing=1.05)
box(s, Inches(0.95), Inches(5.0), Inches(10), Inches(1.0),
    "Cómo el dashboard mide el desempeño del equipo de ventas\na partir de las conversaciones de WhatsApp / Atom",
    size=19, color=SAGEL, line_spacing=1.3)
rect(s, Inches(0.95), Inches(6.55), Inches(2.6), Inches(0.06), SAGE, radius=False)

# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — Qué es el PXI
# ════════════════════════════════════════════════════════════════════
s = slide()
tag(s, Inches(0.9), Inches(0.6), "El concepto")
box(s, Inches(0.9), Inches(1.0), Inches(11), Inches(1.0),
    "Una sola calificación por agente", size=38, color=FOREST, bold=True)
box(s, Inches(0.9), Inches(2.05), Inches(7.0), Inches(2.2),
    "El PXI resume, en un número de 0 a 100, qué tan bien cada agente "
    "cuida la experiencia del paciente durante la conversación de ventas.\n\n"
    "Se calcula automáticamente leyendo el texto de los chats — no requiere "
    "que nadie califique a mano.", size=17, color=FSOFT, line_spacing=1.4)

# big PXI sample card
rect(s, Inches(8.3), Inches(2.0), Inches(4.1), Inches(3.3), IVORY, BORDER, Pt(1))
box(s, Inches(8.5), Inches(2.35), Inches(3.7), Inches(0.5), "EJEMPLO · MARÍA LÓPEZ", size=12, color=FSOFT, bold=True)
box(s, Inches(8.5), Inches(2.7), Inches(3.7), Inches(1.5), "84", size=96, color=OK, bold=True, align=PP_ALIGN.CENTER)
box(s, Inches(8.5), Inches(4.3), Inches(3.7), Inches(0.4), "PXI / 100", size=14, color=FSOFT, bold=True, align=PP_ALIGN.CENTER)
box(s, Inches(8.5), Inches(4.75), Inches(3.7), Inches(0.4), "32 conversaciones · 11 citas", size=12, color=FSOFT, align=PP_ALIGN.CENTER)

box(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.0),
    "El PXI se construye con 5 pilares (auto-evaluados) y se complementa con "
    "alertas de intervención inmediata y un muestreo manual del supervisor.",
    size=15, color=FOREST, italic=True, line_spacing=1.3)

# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — Los 5 pilares + pesos
# ════════════════════════════════════════════════════════════════════
s = slide()
tag(s, Inches(0.9), Inches(0.55), "Cómo se calcula")
box(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.9),
    "Los 5 pilares y su peso", size=36, color=FOREST, bold=True)
box(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(0.6),
    "Cada pilar se evalúa solo cuando aplica a la conversación. El PXI es el promedio "
    "ponderado de los pilares aplicables.", size=15, color=FSOFT, line_spacing=1.3)

pillars = [
    ("P1", "Velocidad", "20%", "¿Respondió al lead en ≤15 minutos?", CDMX),
    ("P2", "Atención plena", "25%", "¿Evitó dejar al paciente esperando respuesta?", SAGE),
    ("P3", "Valor antes de precio", "15%", "¿Acompañó el precio con información del servicio?", WARN),
    ("P4", "Lenguaje seguro", "25%", "¿Evitó promesas clínicas y vocabulario prohibido?", BAD),
    ("P5", "Sensibilidad emocional", "15%", "Si hubo carga emocional, ¿la validó primero?", OK),
]
y = Inches(2.6)
for code, name, wt, desc, c in pillars:
    rect(s, Inches(0.9), y, Inches(11.5), Inches(0.78), IVORY, BORDER, Pt(1))
    rect(s, Inches(0.9), y, Inches(0.14), Inches(0.78), c, radius=False)
    box(s, Inches(1.25), y, Inches(0.9), Inches(0.78), code, size=20, color=c, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(2.2), y, Inches(3.2), Inches(0.78), name, size=17, color=FOREST, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(5.5), y, Inches(5.3), Inches(0.78), desc, size=14, color=FSOFT, anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(10.9), y, Inches(1.3), Inches(0.78), wt, size=22, color=c, bold=True,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    y = y + Inches(0.88)

# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — Cómo se obtiene el número (fórmula visual)
# ════════════════════════════════════════════════════════════════════
s = slide(FOREST)
tag(s, Inches(0.9), Inches(0.6), "La mecánica")
box(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.9),
    "De conversación a calificación", size=36, color=CREAM, bold=True)

steps = [
    ("1", "Agrupar", "Se juntan todos los mensajes de una misma conversación y se ordenan por hora."),
    ("2", "Filtrar", "Solo se auditan conversaciones de ventas reales (con cliente y agente humano)."),
    ("3", "Evaluar", "Cada pilar aplicable se marca CUMPLE o NO CUMPLE según el texto del chat."),
    ("4", "Ponderar", "Se promedian los pilares aplicables con sus pesos → PXI de 0 a 100."),
]
x = Inches(0.9)
for n, t, d in steps:
    rect(s, x, Inches(2.2), Inches(2.85), Inches(3.2), IVORY)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), Inches(2.55), Inches(0.85), Inches(0.85))
    c.fill.solid(); c.fill.fore_color.rgb = SAGE; c.line.fill.background(); c.shadow.inherit = False
    box(s, x + Inches(0.3), Inches(2.55), Inches(0.85), Inches(0.85), n, size=30, color=CREAM, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, x + Inches(0.3), Inches(3.65), Inches(2.3), Inches(0.5), t, size=20, color=FOREST, bold=True)
    box(s, x + Inches(0.3), Inches(4.2), Inches(2.3), Inches(1.1), d, size=13, color=FSOFT, line_spacing=1.3)
    x = x + Inches(2.95)

box(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.9),
    "Un agente nunca es penalizado por un pilar que no aplicó en su conversación: "
    "el promedio se re-normaliza solo sobre los pilares evaluables.",
    size=15, color=SAGEL, italic=True, line_spacing=1.35)

# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — Alertas R1–R4
# ════════════════════════════════════════════════════════════════════
s = slide()
tag(s, Inches(0.9), Inches(0.55), "Intervención inmediata")
box(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.9),
    "4 alertas que se escalan el mismo día", size=34, color=FOREST, bold=True)
box(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(0.6),
    "Son señales binarias (ocurrió / no ocurrió). NO afectan el PXI: van a una cola aparte "
    "con enlace directo al chat en Atom.", size=15, color=FSOFT, line_spacing=1.3)

flags = [
    ("R1", "Frustración no validada", "El paciente se mostró molesto y no hubo una respuesta empática.", BAD),
    ("R2", "Info clínica sin confirmación", "Se dieron tasas de éxito, garantías o diagnósticos sin aval médico.", BAD),
    ("R3", "Alto valor sin supervisor", "Se cotizó FIV, PGT u ovodonación; conviene revisar apoyo de supervisora.", CDMX),
    ("R4", "Lead >24h sin respuesta", "Un lead interesado quedó sin contestación por más de 24 horas.", SAGE),
]
positions = [(Inches(0.9), Inches(2.6)), (Inches(6.75), Inches(2.6)),
             (Inches(0.9), Inches(4.55)), (Inches(6.75), Inches(4.55))]
for (code, name, desc, c), (x, y) in zip(flags, positions):
    rect(s, x, y, Inches(5.65), Inches(1.7), IVORY, BORDER, Pt(1))
    rect(s, x, y, Inches(0.14), Inches(1.7), c, radius=False)
    bdg = rect(s, x + Inches(0.35), y + Inches(0.32), Inches(0.95), Inches(0.55), c)
    box(s, x + Inches(0.35), y + Inches(0.32), Inches(0.95), Inches(0.55), code, size=20, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, x + Inches(1.5), y + Inches(0.28), Inches(3.9), Inches(0.5), name, size=16, color=FOREST, bold=True)
    box(s, x + Inches(1.5), y + Inches(0.78), Inches(3.95), Inches(0.8), desc, size=12.5, color=FSOFT, line_spacing=1.25)

# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — Muestreo manual + limitación
# ════════════════════════════════════════════════════════════════════
s = slide()
tag(s, Inches(0.9), Inches(0.55), "Lo que el texto no puede medir")
box(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.9),
    "Muestreo manual del supervisor", size=34, color=FOREST, bold=True)

rect(s, Inches(0.9), Inches(2.0), Inches(6.5), Inches(3.9), IVORY, BORDER, Pt(1))
box(s, Inches(1.2), Inches(2.25), Inches(6), Inches(0.5), "El supervisor revisa ~3 chats por agente y marca:", size=15, color=FOREST, bold=True)
items = [
    "Uso real del nombre / personalización genuina",
    "≥2 preguntas abiertas (SPIN) antes de informar",
    "Protocolo de objeciones completo",
    "Siguiente paso concreto y fechado al cierre",
    "Motivo de cita y estado del deal en HubSpot",
]
yy = Inches(2.85)
for it in items:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), yy + Inches(0.05), Inches(0.22), Inches(0.22))
    c.fill.solid(); c.fill.fore_color.rgb = SAGE; c.line.fill.background(); c.shadow.inherit = False
    box(s, Inches(1.6), yy, Inches(5.6), Inches(0.55), it, size=14.5, color=FSOFT, line_spacing=1.2)
    yy = yy + Inches(0.58)

rect(s, Inches(7.7), Inches(2.0), Inches(4.7), Inches(3.9), FOREST)
box(s, Inches(8.0), Inches(2.3), Inches(4.1), Inches(0.5), "⚠  LIMITACIÓN CLAVE", size=14, color=WARN, bold=True)
box(s, Inches(8.0), Inches(2.85), Inches(4.1), Inches(2.9),
    "El export de Atom es solo texto.\n\n"
    "Los precios, paquetes e información que se envían como imagen o PDF no aparecen "
    "en el contenido.\n\n"
    "Por eso solo se auto-evalúa lo verificable por texto; el resto lo revisa una persona.",
    size=14.5, color=CREAM, line_spacing=1.35)

# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — Cómo leer el dashboard
# ════════════════════════════════════════════════════════════════════
s = slide(FOREST)
tag(s, Inches(0.9), Inches(0.6), "En el día a día")
box(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.9),
    "Cómo usar el dashboard", size=36, color=CREAM, bold=True)

tabs = [
    ("Scorecard PXI", "El PXI de cada agente y sus 5 pilares con colores. Verde ≥80, amarillo 60–79, rojo <60. Identifica de un vistazo qué pilar arrastra a cada quien."),
    ("Intervención inmediata", "La cola de alertas R1–R4, filtrable por agente y tipo, con enlace directo a la conversación en Atom para actuar el mismo día."),
    ("Muestreo manual", "El checklist que el supervisor llena sobre los chats muestreados. Las marcas se guardan en el navegador."),
]
y = Inches(2.3)
for t, d in tabs:
    rect(s, Inches(0.9), y, Inches(11.5), Inches(1.35), IVORY)
    box(s, Inches(1.25), y + Inches(0.2), Inches(3.4), Inches(1.0), t, size=19, color=SAGE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    box(s, Inches(4.8), y + Inches(0.2), Inches(7.3), Inches(1.0), d, size=14, color=FSOFT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
    y = y + Inches(1.5)

box(s, Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.4),
    "Actualizar datos: sube el nuevo export como  datos.xlsx  en GitHub → Add file → Upload files.",
    size=13, color=SAGEL, italic=True)

prs.save("PXI_Como_mide_el_desempeno.pptx")
print("Guardado: PXI_Como_mide_el_desempeno.pptx ·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
